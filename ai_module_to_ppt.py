import os
import json
import argparse
from docx import Document
from pptx import Presentation
from jinja2 import Template
from openai import AzureOpenAI
from dotenv import load_dotenv
import prompty
from pathlib import Path
import logging

load_dotenv()

client = AzureOpenAI(api_key=os.environ['AZURE_OPENAI_API_KEY'],
azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
api_version="2025-01-01-preview")

# Configure logging
def setup_logging():
    logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# 1) Extract document sections
def extract_sections(docx_path):
    sections = []
    doc = Document(docx_path)
    current = None
    for para in doc.paragraphs:
        txt = para.text.strip()
        if not txt:
            continue
        style = para.style.name.lower()
        if 'heading' in style:
            if current:
                sections.append(current)
            current = { 'title': txt, 'content': [] }
        elif current:
            current['content'].append(txt)
    if current:
        sections.append(current)
    logging.info(f"Extracted {len(sections)} sections from {docx_path}")
    return sections

# 2) Use prompty to generate slide specs
def generate_slides_ai(sections, prompty_path):
    """Generate slide specifications using the prompty library"""
    # We'll use Azure OpenAI directly for simplicity to avoid encoding issues with prompty
    slides = []
    
    try:
        # Prepare a system prompt and user message
        system_prompt = """
        You're an AI presentation architect. Convert document sections into concise, engaging PowerPoint slides.
        For each section:
        1. Create a clear title slide
        2. Extract 4-6 key points as bullet points
        3. Ensure consistency of tone and style
        4. Focus on impactful information

        The response format must be a valid JSON array of slide objects:
        [
            {
                "title": "Slide Title",
                "bullets": ["Bullet point 1", "Bullet point 2", "..."]
            }
        ]

        Include a title slide and a closing/questions slide as bookends.
        """
        
        user_prompt = f"""
        I need to convert these document sections into PowerPoint slides.
        
        Document sections:
        {json.dumps(sections, ensure_ascii=False)}
        
        Please generate a presentation structure with appropriate slides, following best practices for presentations.
        """
        
        # Call Azure OpenAI
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"}
        )
        
        raw = response.choices[0].message.content
        result = json.loads(raw)
        
        # Extract slides from the result
        if isinstance(result, dict) and "slides" in result:
            slides = result["slides"]
        elif isinstance(result, list):
            slides = result
            
        logging.info(f"Generated {len(slides)} slides via AI")
        return slides
        
    except Exception as e:
        logging.error(f"Error generating slides: {str(e)}")
        # Return a basic set of slides based on sections
        for section in sections:
            if section['title'] not in ['Feedback']:  # Skip feedback sections
                bullet_content = section.get('content', [])
                bullets = [item for item in bullet_content if not item.startswith('•') and len(item) > 10][:5]
                slides.append({
                    "title": section['title'],
                    "bullets": bullets if bullets else ["Content to be added"]
                })
        
        logging.info(f"Generated {len(slides)} basic slides after error")
        return slides

# 4) Apply template and build PPT
def build_presentation(template_path, slide_specs, output_path):
    prs = Presentation(template_path)
    # remove default slides
    for _ in range(len(prs.slides)):
        rId = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(rId)
    # add slides per spec
    for spec in slide_specs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = spec.get('title', '')
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for b in spec.get('bullets', []):
            p = body.add_paragraph(); p.text = b; p.level = 0
    prs.save(output_path)
    logging.info(f"Saved AI-enhanced presentation: {output_path}")

# CLI entrypoint
def main():
    setup_logging()
    p = argparse.ArgumentParser(description='AI-driven Word→PPT using template')
    p.add_argument('-d','--docx', metavar='DOCX', default='develop-ai-agent-with-semantic-kernel.docx', help='Word file to convert (default: develop-ai-agent-with-semantic-kernel.docx)')
    p.add_argument('-t','--template', metavar='PPTX', default='speaker-template.pptx', help='PPTX template to use (default: speaker-template.pptx)')
    p.add_argument('-o','--output', default='output_presentation.pptx', help='Output PPTX file')
    p.add_argument('--prompty', default='slide_generator.prompty', help='Prompty file for slide AI')
    args = p.parse_args()

    secs = extract_sections(args.docx)
    specs = generate_slides_ai(secs, args.prompty)
    build_presentation(args.template, specs, args.output)

if __name__=='__main__':
    main()

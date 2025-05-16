# Use only essential imports
import os
import sys
import argparse
import logging
from docx import Document
from pptx import Presentation

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class WordToPowerPointConverter:
    """Vanilla converter: converts Word sections into PPT slides"""
    def __init__(self):
        pass

    def extract_document_structure(self, docx_path):
        """Extract headings and paragraphs from Word document"""
        doc = Document(docx_path)
        sections = []
        current = None
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith('Heading'):
                if current:
                    sections.append(current)
                current = {"title": text, "content": []}
            elif current:
                current["content"].append(text)
        if current:
            sections.append(current)
        return sections

    def convert_document(self, docx_path, output_pptx_path):
        """Convert Word document to a basic PowerPoint deck"""
        logger.info(f"Converting {docx_path} to {output_pptx_path}")
        prs = Presentation()
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = os.path.splitext(os.path.basename(docx_path))[0]
        # extract sections
        sections = self.extract_document_structure(docx_path)
        for sec in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec['title']
            if sec['content']:
                tf = slide.shapes.placeholders[1].text_frame
                tf.clear()
                for para in sec['content']:
                    p = tf.add_paragraph()
                    p.text = para
                    p.level = 0
        prs.save(output_pptx_path)
        logger.info(f"Saved presentation: {output_pptx_path}")
        return output_pptx_path

def main():
    parser = argparse.ArgumentParser(description='Convert Word document to PowerPoint')
    parser.add_argument('input_file', help='Path to the Word document')
    parser.add_argument('--output', '-o', help='Path to save the PowerPoint file')
    args = parser.parse_args()

    input_path = args.input_file
    output_path = args.output or os.path.splitext(input_path)[0] + '_presentation.pptx'
    
    try:
        converter = WordToPowerPointConverter()
        converter.convert_document(input_path, output_path)
        print(f"✅ Created presentation: {output_path}")
    except Exception as err:
        print(f"❌ Error: {err}")
        sys.exit(1)

if __name__ == "__main__":
    main()
